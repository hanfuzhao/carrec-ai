"""Build the CarRec AI pitch deck as a PPTX file.

Editorial automotive-magazine style: ink background, paper text,
rust/brass/forest accents. Slides match pitch_presentation.html.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


INK = RGBColor(0x1A, 0x16, 0x12)
PAPER = RGBColor(0xEB, 0xE5, 0xD8)
RUST = RGBColor(0xC8, 0x36, 0x2D)
BRASS = RGBColor(0xC4, 0xA3, 0x5A)
FOREST = RGBColor(0x2D, 0x4A, 0x2B)
DIM = RGBColor(0x8A, 0x80, 0x70)

DISPLAY = "Georgia"
BODY = "Georgia"
MONO = "Courier New"


def add_slide(prs, bg=INK):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg
    bg_shape.line.fill.background()
    bg_shape.shadow.inherit = False
    bg_shape._element.addprevious(bg_shape._element)
    return slide


def add_text(slide, x, y, w, h, text, *, font=BODY, size=14, color=PAPER,
             bold=False, italic=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.15, tracking=0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        if tracking:
            rPr = run._r.get_or_add_rPr()
            rPr.set("spc", str(tracking))
    return tb


def add_rect(slide, x, y, w, h, color, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_line(slide, x, y, w, color, weight=1.0):
    ln = slide.shapes.add_connector(1, x, y, x + w, y)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def set_notes(slide, text):
    notes = slide.notes_slide
    notes.notes_text_frame.text = text


def slide_header(slide, num, section, prs):
    add_text(slide, Inches(0.6), Inches(0.45), Inches(6), Inches(0.3),
             f"{num} / {section}", font=MONO, size=10, color=DIM, tracking=200)
    add_text(slide, Inches(11.0), Inches(0.45), Inches(1.7), Inches(0.3),
             "CarRec AI", font=MONO, size=10, color=DIM, align=PP_ALIGN.RIGHT,
             tracking=200)
    add_line(slide, Inches(0.6), Inches(0.85), Inches(12.13), DIM, 0.5)


def slide_footer(slide, page, total):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(6), Inches(0.3),
             "Module 3 Hackathon / Recommendation Systems",
             font=MONO, size=8, color=DIM, tracking=100)
    add_text(slide, Inches(11.0), Inches(7.05), Inches(1.7), Inches(0.3),
             f"{page:02d} / {total:02d}", font=MONO, size=8, color=DIM,
             align=PP_ALIGN.RIGHT, tracking=100)


def build_slide_1(prs):
    s = add_slide(prs)
    add_text(s, Inches(0.6), Inches(0.45), Inches(8), Inches(0.3),
             "01 / Title", font=MONO, size=10, color=DIM, tracking=200)
    add_text(s, Inches(11.0), Inches(0.45), Inches(1.7), Inches(0.3),
             "CarRec AI", font=MONO, size=10, color=DIM, align=PP_ALIGN.RIGHT,
             tracking=200)
    add_line(s, Inches(0.6), Inches(0.85), Inches(12.13), DIM, 0.5)

    add_text(s, Inches(0.6), Inches(1.4), Inches(8), Inches(0.3),
             "PITCH 2026 / CARREC AI", font=MONO, size=11, color=BRASS,
             tracking=400)

    add_text(s, Inches(0.5), Inches(2.1), Inches(9), Inches(2.2),
             "CarRec", font=DISPLAY, size=110, color=PAPER, bold=True,
             line_spacing=0.95)
    add_text(s, Inches(0.5), Inches(3.6), Inches(9), Inches(2.2),
             "AI.", font=DISPLAY, size=110, color=PAPER, bold=True,
             line_spacing=0.95)

    add_rect(s, Inches(0.6), Inches(5.55), Inches(0.5), Inches(0.06), RUST)

    add_text(s, Inches(0.6), Inches(5.75), Inches(8), Inches(0.9),
             "Recommendations that actually matter. An LLM-powered car\n"
             "recommendation system balancing relevance with responsibility.",
             font=BODY, size=15, color=PAPER, italic=True, line_spacing=1.3)

    add_text(s, Inches(0.6), Inches(6.85), Inches(6), Inches(0.3),
             "github.com/hanfuzhao/carrec-ai", font=MONO, size=10, color=BRASS)
    add_text(s, Inches(6.8), Inches(6.85), Inches(6), Inches(0.3),
             "carrec-ai.onrender.com", font=MONO, size=10, color=BRASS,
             align=PP_ALIGN.RIGHT)

    set_notes(s, "This is CarRec AI, a car recommendation system that uses an "
                 "LLM to balance relevance with responsibility. The live app "
                 "and source code are linked on this slide.")


def build_slide_2(prs):
    s = add_slide(prs)
    slide_header(s, "02", "The Problem", prs)

    add_text(s, Inches(0.6), Inches(1.25), Inches(7.2), Inches(0.4),
             "THE PROBLEM", font=MONO, size=11, color=RUST, tracking=400)

    add_text(s, Inches(0.6), Inches(1.7), Inches(7.2), Inches(2.6),
             "Engagement optimization\ncreates feedback loops\n"
             "that narrow user\nchoice.",
             font=DISPLAY, size=44, color=PAPER, bold=True, line_spacing=1.05)

    add_text(s, Inches(0.6), Inches(4.7), Inches(6.8), Inches(2),
             "Car platforms like Dongchedi and Edmunds push popular\n"
             "models because popular models get clicks. But a first-time\n"
             "buyer asking for a family SUV does not need to see five\n"
             "Teslas.",
             font=BODY, size=13, color=DIM, italic=True, line_spacing=1.35)

    add_line(s, Inches(8.1), Inches(1.25), Inches(0.02), DIM, 0.5)

    items = [
        ("01", "Popular cars", "get more exposure, become more popular, "
                               "crowd out alternatives."),
        ("02", "Niche brands", "(BYD, NIO, Rivian) get buried even when "
                                "they fit better."),
        ("03", "Users see", "five Teslas when they ask for a family SUV."),
        ("04", "No budget enforcement,", "no diversity, no fairness."),
    ]
    y = 1.4
    for num, head, body in items:
        add_text(s, Inches(8.4), Inches(y), Inches(0.6), Inches(0.3),
                 num, font=MONO, size=11, color=BRASS, tracking=200)
        add_text(s, Inches(9.0), Inches(y), Inches(3.7), Inches(1.0),
                 head + " " + body, font=BODY, size=13, color=PAPER,
                 line_spacing=1.3)
        y += 1.3

    slide_footer(s, 2, 8)
    set_notes(s, "Engagement-optimized recommenders create feedback loops. "
                 "Popular cars get more exposure, become more popular, and "
                 "crowd out alternatives. A first-time buyer asking for a "
                 "family SUV does not need to see five Teslas. They need "
                 "options that fit their budget and expose them to brands "
                 "they might not have considered.")


def build_slide_3(prs):
    s = add_slide(prs)
    slide_header(s, "03", "What I Built", prs)

    add_text(s, Inches(0.6), Inches(1.25), Inches(8), Inches(1.6),
             "What\nI Built.", font=DISPLAY, size=64, color=PAPER, bold=True,
             line_spacing=0.95)

    add_text(s, Inches(0.6), Inches(3.55), Inches(7.5), Inches(1),
             "GPT-4o-mini with retrieval-augmented generation,\n"
             "enforcing four real-world constraints.",
             font=BODY, size=14, color=DIM, italic=True, line_spacing=1.3)

    cards = [
        ("01", "Budget\nCompliance", "Hard filter. Never recommends above "
                                      "the user's stated budget."),
        ("02", "Brand\nDiversity", "Greedy selection penalizing repeated "
                                    "brands to ensure variety."),
        ("03", "Fairness\nBoost", "Scoring advantage for niche brands to "
                                   "ensure equitable exposure."),
        ("04", "Cold\nStart", "Preference parsing from natural language. "
                              "Zero user history needed."),
    ]
    x = 0.6
    for num, title, desc in cards:
        add_rect(s, Inches(x), Inches(4.8), Inches(2.95), Inches(2.0),
                 RGBColor(0x22, 0x1D, 0x18))
        add_text(s, Inches(x + 0.25), Inches(4.95), Inches(2.5), Inches(0.3),
                 num, font=MONO, size=11, color=BRASS, tracking=200)
        add_text(s, Inches(x + 0.25), Inches(5.3), Inches(2.5), Inches(0.9),
                 title, font=DISPLAY, size=20, color=PAPER, bold=True,
                 line_spacing=1.0)
        add_text(s, Inches(x + 0.25), Inches(6.2), Inches(2.5), Inches(0.6),
                 desc, font=BODY, size=10, color=DIM, line_spacing=1.25)
        x += 3.1

    slide_footer(s, 3, 8)
    set_notes(s, "I built a car recommendation system using GPT-4o-mini from "
                 "OpenAI, adapted to the car domain through prompt engineering "
                 "and retrieval-augmented generation. The system enforces four "
                 "real-world constraints: budget compliance as a hard filter, "
                 "brand diversity through greedy selection, fairness via a "
                 "scoring boost for niche brands, and cold start handling by "
                 "parsing preferences from the query itself with zero user "
                 "history.")


def build_slide_4(prs):
    s = add_slide(prs)
    slide_header(s, "04", "LLM Strategy", prs)

    add_text(s, Inches(0.6), Inches(1.2), Inches(8), Inches(0.4),
             "04 / LLM STRATEGY", font=MONO, size=11, color=RUST, tracking=400)

    add_text(s, Inches(0.6), Inches(1.65), Inches(12), Inches(1.2),
             "Three layers of adaptation.", font=DISPLAY, size=44,
             color=PAPER, bold=True, line_spacing=1.0)

    layers = [
        ("01", "Structured Retrieval",
         "Parse query to extract budget, use cases, energy, body, seats. "
         "Filter catalog of 70 cars across 27 brands. Only relevant "
         "candidates reach the LLM."),
        ("02", "Constraint-Aware Scoring",
         "Each car scored on use-case match, preference alignment, budget "
         "utilization, and fairness boost. Diversity-aware selection ensures "
         "top 5 span multiple brands."),
        ("03", "LLM Reasoning",
         "GPT-4o-mini generates personalized 1-2 sentence reasons for each "
         "recommendation. Falls back to templates without API key. "
         "Correctness never depends on the LLM."),
    ]
    y = 3.0
    for num, name, desc in layers:
        add_text(s, Inches(0.6), Inches(y), Inches(0.8), Inches(0.4),
                 num, font=MONO, size=14, color=BRASS, tracking=300)
        add_line(s, Inches(1.5), Inches(y + 0.18), Inches(0.5), BRASS, 1.5)
        add_text(s, Inches(2.2), Inches(y), Inches(3.5), Inches(0.5),
                 name, font=DISPLAY, size=20, color=PAPER, bold=True)
        add_text(s, Inches(6.0), Inches(y - 0.05), Inches(6.8), Inches(1.2),
                 desc, font=BODY, size=12, color=DIM, line_spacing=1.3)
        y += 1.25

    add_text(s, Inches(0.6), Inches(6.85), Inches(12), Inches(0.3),
             "Adapted via prompt engineering + RAG. No weight modification.",
             font=BODY, size=11, color=BRASS, italic=True)

    slide_footer(s, 4, 8)
    set_notes(s, "The strategy has three layers. First, structured retrieval. "
                 "The LLM never sees the raw query alone. The system parses "
                 "it to extract budget, use cases, and preferences, then "
                 "filters a catalog of 70 cars. Only relevant candidates move "
                 "forward. Second, constraint-aware scoring. Each car is "
                 "scored on use-case match, preference alignment, budget "
                 "utilization, and a fairness boost for niche brands. A "
                 "diversity-aware selection ensures variety. Third, LLM "
                 "reasoning. GPT-4o-mini generates a personalized reason for "
                 "each recommendation. When no API key is available, the "
                 "system falls back to templates and still produces correct "
                 "recommendations. The LLM enhances the experience but is "
                 "not a dependency for correctness.")


def build_slide_5(prs):
    s = add_slide(prs)
    slide_header(s, "05", "Evaluation", prs)

    add_text(s, Inches(0.6), Inches(1.2), Inches(8), Inches(0.4),
             "05 / EVALUATION", font=MONO, size=11, color=RUST, tracking=400)

    add_text(s, Inches(0.6), Inches(1.65), Inches(8), Inches(1.2),
             "Before vs. After.", font=DISPLAY, size=48, color=PAPER,
             bold=True, line_spacing=1.0)

    add_text(s, Inches(0.6), Inches(2.8), Inches(12), Inches(0.5),
             "15 test queries covering family, luxury, electric, pickup, "
             "and budget use cases.",
             font=BODY, size=13, color=DIM, italic=True)

    table_y = 3.5
    row_h = 0.55
    cols = [
        ("Metric", 4.6),
        ("Naive", 2.3),
        ("Smart", 2.3),
        ("Change", 2.93),
    ]
    x = 0.6
    for label, w in cols:
        add_text(s, Inches(x), Inches(table_y), Inches(w), Inches(0.4),
                 label.upper(), font=MONO, size=10, color=DIM, tracking=200)
        x += w
    add_line(s, Inches(0.6), Inches(table_y + 0.4), Inches(12.13), PAPER, 1.0)

    rows = [
        ("Budget compliance", "84%", "100%", "+19%", RUST),
        ("Brand diversity", "3.0", "4.93", "+64%", FOREST),
        ("Type diversity", "2.0", "2.53", "+27%", FOREST),
        ("Niche exposure", "0%", "60%", "from zero", BRASS),
        ("Relevance", "0.38", "0.79", "+109%", RUST),
    ]
    y = table_y + 0.6
    for metric, naive, smart, change, change_color in rows:
        x = 0.6
        add_text(s, Inches(x), Inches(y), Inches(4.6), Inches(row_h),
                 metric, font=BODY, size=14, color=PAPER)
        x += 4.6
        add_text(s, Inches(x), Inches(y), Inches(2.3), Inches(row_h),
                 naive, font=MONO, size=14, color=DIM)
        x += 2.3
        add_text(s, Inches(x), Inches(y), Inches(2.3), Inches(row_h),
                 smart, font=MONO, size=16, color=PAPER, bold=True)
        x += 2.3
        add_text(s, Inches(x), Inches(y), Inches(2.93), Inches(row_h),
                 change, font=MONO, size=14, color=change_color, bold=True)
        y += row_h
        add_line(s, Inches(0.6), Inches(y - 0.05), Inches(12.13), DIM, 0.25)

    slide_footer(s, 5, 8)
    set_notes(s, "I evaluated both modes across 15 test queries. The naive "
                 "mode, which sorts by popularity, achieved 84 percent budget "
                 "compliance, meaning 16 percent of recommendations exceeded "
                 "the user's budget. Brand diversity was stuck at 3, and "
                 "niche brand exposure was zero. The smart mode achieved 100 "
                 "percent budget compliance, increased brand diversity to "
                 "nearly 5 unique brands per query, and brought niche "
                 "exposure from zero to 60 percent. Relevance more than "
                 "doubled from 0.38 to 0.79.")


def build_slide_6(prs):
    s = add_slide(prs)
    slide_header(s, "06", "Example Query", prs)

    add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.4),
             "06 / EXAMPLE QUERY", font=MONO, size=11, color=RUST, tracking=400)

    add_text(s, Inches(0.6), Inches(1.8), Inches(12), Inches(1),
             "\u201cI need a family SUV under $50k, preferably electric\u201d",
             font=DISPLAY, size=26, color=PAPER, italic=True, bold=True,
             line_spacing=1.15)

    col_y = 3.1
    col_h = 3.6

    add_rect(s, Inches(0.6), Inches(col_y), Inches(5.9), Inches(col_h),
             RGBColor(0x22, 0x1D, 0x18))
    add_text(s, Inches(0.85), Inches(col_y + 0.2), Inches(5.5), Inches(0.4),
             "BEFORE / NAIVE", font=MONO, size=11, color=DIM, tracking=300)
    add_line(s, Inches(0.85), Inches(col_y + 0.65), Inches(5.4), DIM, 0.5)

    before_cars = [
        ("Tesla Model Y", "$54k"),
        ("Tesla Model 3", "$52k"),
        ("Toyota RAV4", "$38k"),
        ("Ford F-150", "$45k"),
        ("Honda CR-V", "$36k"),
    ]
    cy = col_y + 0.85
    for name, price in before_cars:
        add_text(s, Inches(0.85), Inches(cy), Inches(4.0), Inches(0.35),
                 name, font=BODY, size=13, color=PAPER)
        add_text(s, Inches(4.85), Inches(cy), Inches(1.4), Inches(0.35),
                 price, font=MONO, size=12, color=DIM, align=PP_ALIGN.RIGHT)
        cy += 0.4
    add_text(s, Inches(0.85), Inches(col_y + col_h - 0.55), Inches(5.5),
             Inches(0.5),
             "2 over budget. All mainstream. 0 electric SUVs. 0 niche brands.",
             font=BODY, size=11, color=RUST, italic=True, line_spacing=1.25)

    add_rect(s, Inches(6.85), Inches(col_y), Inches(5.9), Inches(col_h),
             RGBColor(0x2A, 0x20, 0x14))
    add_text(s, Inches(7.1), Inches(col_y + 0.2), Inches(5.5), Inches(0.4),
             "AFTER / SMART", font=MONO, size=11, color=BRASS, tracking=300)
    add_line(s, Inches(7.1), Inches(col_y + 0.65), Inches(5.4), BRASS, 0.5)

    after_cars = [
        ("BYD Atto 3", "$28k"),
        ("NIO ES6", "$42k"),
        ("Volvo XC40 Recharge", "$42k"),
        ("VW ID.4", "$39k"),
        ("Hyundai Ioniq 5", "$42k"),
    ]
    cy = col_y + 0.85
    for name, price in after_cars:
        add_text(s, Inches(7.1), Inches(cy), Inches(4.0), Inches(0.35),
                 name, font=BODY, size=13, color=PAPER)
        add_text(s, Inches(11.1), Inches(cy), Inches(1.4), Inches(0.35),
                 price, font=MONO, size=12, color=BRASS, align=PP_ALIGN.RIGHT)
        cy += 0.4
    add_text(s, Inches(7.1), Inches(col_y + col_h - 0.55), Inches(5.5),
             Inches(0.5),
             "All within budget. All electric SUVs. 3 niche brands exposed.",
             font=BODY, size=11, color=FOREST, italic=True, line_spacing=1.25)

    slide_footer(s, 6, 8)
    set_notes(s, "For the query 'family SUV under 50k, preferably electric,' "
                 "the naive mode returns the five most popular cars overall. "
                 "Two exceed the budget, none are electric SUVs, and all are "
                 "mainstream brands. The smart mode returns five electric "
                 "SUVs, all within budget, from three niche brands that the "
                 "user might not have considered otherwise.")


def build_slide_7(prs):
    s = add_slide(prs)
    slide_header(s, "07", "Risks & Ethics", prs)

    add_text(s, Inches(0.6), Inches(1.25), Inches(8), Inches(1.2),
             "Risks & Ethics.", font=DISPLAY, size=48, color=PAPER,
             bold=True, line_spacing=1.0)
    add_text(s, Inches(0.6), Inches(2.35), Inches(8), Inches(0.4),
             "07 / REFLECTION", font=MONO, size=11, color=RUST, tracking=400)

    risks = [
        ("01", "Catalog Bias",
         "70 cars is curated, not exhaustive. Niche classification is manual "
         "and should be data-driven in a production system."),
        ("02", "Fairness Is Editorial",
         "The 1.5-point niche boost is a value judgment, not objective "
         "optimization. Different stakeholders may disagree on the right "
         "amount."),
        ("03", "No Learning Loop",
         "Cold start means no feedback mechanism. The system cannot improve "
         "from outcomes without tracking whether users acted on "
         "recommendations."),
        ("04", "LLM Hallucination",
         "The LLM generates reasons but does not make the decision. "
         "Structured data validates the recommendation. Production would "
         "validate LLM output against car data."),
    ]
    positions = [
        (0.6, 3.0), (6.85, 3.0), (0.6, 5.05), (6.85, 5.05)
    ]
    for (num, title, desc), (x, y) in zip(risks, positions):
        add_text(s, Inches(x), Inches(y), Inches(0.6), Inches(0.4),
                 num, font=MONO, size=12, color=BRASS, tracking=300)
        add_line(s, Inches(x + 0.5), Inches(y + 0.18), Inches(0.4), BRASS, 1.0)
        add_text(s, Inches(x + 1.0), Inches(y), Inches(4.8), Inches(0.5),
                 title, font=DISPLAY, size=18, color=PAPER, bold=True)
        add_text(s, Inches(x + 1.0), Inches(y + 0.45), Inches(4.8), Inches(1.5),
                 desc, font=BODY, size=11, color=DIM, line_spacing=1.3)

    slide_footer(s, 7, 8)
    set_notes(s, "Four risks to consider. First, catalog bias. My 70-car "
                 "catalog is curated, and I chose which brands count as "
                 "niche. A real system would need a larger, continuously "
                 "updated catalog with data-driven classification. Second, "
                 "fairness is not neutral. The 1.5-point boost for niche "
                 "brands is an editorial choice. Third, there is no learning "
                 "loop. Every session is a cold start with no feedback "
                 "mechanism. Fourth, LLM hallucination. The LLM generates "
                 "reasons but does not make the recommendation decision, "
                 "which is deterministic and auditable. But the reasons "
                 "could mislead users if the LLM invents attributes. A "
                 "production system would validate LLM output against "
                 "structured data.")


def build_slide_8(prs):
    s = add_slide(prs)
    slide_header(s, "08", "Live Demo", prs)

    add_text(s, Inches(0.6), Inches(1.5), Inches(8), Inches(0.4),
             "08 / DEMO", font=MONO, size=11, color=RUST, tracking=400)

    add_text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(2.5),
             "Try It\nLive.", font=DISPLAY, size=110, color=PAPER, bold=True,
             line_spacing=0.95)

    add_text(s, Inches(0.6), Inches(4.9), Inches(12), Inches(0.5),
             "The app is deployed and accepts natural language queries.",
             font=BODY, size=15, color=DIM, italic=True)

    add_line(s, Inches(0.6), Inches(5.6), Inches(12.13), DIM, 0.5)

    add_text(s, Inches(0.6), Inches(5.8), Inches(3), Inches(0.3),
             "LIVE APP", font=MONO, size=10, color=BRASS, tracking=300)
    add_text(s, Inches(0.6), Inches(6.05), Inches(7), Inches(0.4),
             "carrec-ai.onrender.com", font=MONO, size=14, color=PAPER)

    add_text(s, Inches(7.0), Inches(5.8), Inches(3), Inches(0.3),
             "SOURCE CODE", font=MONO, size=10, color=BRASS, tracking=300)
    add_text(s, Inches(7.0), Inches(6.05), Inches(6), Inches(0.4),
             "github.com/hanfuzhao/carrec-ai", font=MONO, size=14, color=PAPER)

    add_text(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
             "\u201cfamily SUV under $50k, preferably electric\u201d  /  "
             "\u201ccheap commute car around $25k\u201d  /  "
             "\u201cluxury sedan for business, budget $60k\u201d",
             font=BODY, size=11, color=DIM, italic=True)

    slide_footer(s, 8, 8)
    set_notes(s, "The live app is deployed on Render and accepts natural "
                 "language queries. The GitHub repository contains the full "
                 "codebase with a proper branch and PR workflow.")


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    build_slide_4(prs)
    build_slide_5(prs)
    build_slide_6(prs)
    build_slide_7(prs)
    build_slide_8(prs)

    out = "pitch_presentation.pptx"
    prs.save(out)
    print(f"Saved {out} with {len(prs.slides)} slides")


if __name__ == "__main__":
    build()
